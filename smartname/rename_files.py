"""
SmartName is an AI File Renamer that analyzes user given files and renames them.
The user may provide a casing style they prefer to have their filesnames follow.
A guide to casing can be found in CASING_GUIDE.md located in the GitHub repository.
View documentation on GitHub for more info.
https://github.com/mcneilkimberly/SmartName---AI-File-Renamer/tree/main/smartname
"""

import argparse
import os
import tempfile # Helps with extracted text
import base64
import requests
import fitz # Helps with PDFs
import json
from pypdf import PdfReader
from PIL import Image # Helps with images
from docx import Document # Helps with Word Documents
from pptx import Presentation # Helps with PowerPoint Presentations

# CLI - Command-Line Interface
def main():
    """
    Command-line interface for smart file renaming using Ollama AI models
      Dry run (preview): python rename_files.py data/
      Execute renaming: python rename_files.py data/ --execute
      Custom model: python rename_files.py data/ --model modelname
    """
    # Create argument parser object
    parser = argparse.ArgumentParser(description='Smart file renaming tool using Ollama AI models')
    # Add directory argument
    parser.add_argument('directory', type=str, help='Directory containing files to rename')
    # Add execute flag, defaults to False when --execute is not included (for dry-run mode)
    parser.add_argument('--execute', action='store_true', help='Execute the renaming. If not set, performs a dry run')
    # Add model selection argument
    parser.add_argument('--model', type=str, default='llama3.2-vision', help='Specify the Ollama model to use (default: llama3.2-vision)')
    # Add casing style argument
    parser.add_argument('--case', type=str, default='snake', choices=['snake', 'kebab', 'camel', 'pascal', 'lower', 'title'], help='Casing style for new filenames (default: snake)')

    # Parse arguments
    args = parser.parse_args()

    # Does the directory exist?
    if not os.path.isdir(args.directory):
        print(f"Error: Directory '{args.directory}' does not exist")
        return

    # Print mode (dry run or execute)
    mode = "EXECUTE" if args.execute else "DRY RUN"
    print(f"\nMode: {mode}")
    print(f"Using model: {args.model}")
    print(f"Casing style: {args.case}")
    print(f"Directory: {os.path.abspath(args.directory)}\n")

    # Process all files in the directory
    for filename in os.listdir(args.directory):
        file_path = os.path.join(args.directory, filename)
        # Skip if it's not a file
        if not os.path.isfile(file_path): continue
        # Get the file extension
        name, ext = os.path.splitext(filename)
        # Generate new filename
        suggested_name = generate_filename_for_file(file_path, args.model)        
        # If a suggestion was made, proceed to rename
        if suggested_name:
            # Apply casing and sanitization
            suggested_name = sanitize_filename(suggested_name)
            suggested_name = apply_casing(suggested_name, args.case)
            # Add back the extension
            new_filename = f"{suggested_name}{ext}"
            new_file_path = os.path.join(args.directory, new_filename)
            # Print the proposed change
            print(f"'{filename}' -> '{new_filename}'")
            # Execute the rename if in execute mode
            if args.execute:
                try:
                    os.rename(file_path, new_file_path)
                    print("Renamed successfully!")
                except Exception as e:
                    print(f"Error renaming: {e}")
        # Otherwise, skip file, as no suggestion was made due to unsupported type or error
        else:
            print(f"Skipping '{filename}' - unsupported file type or error in processing")
    # Print complete
    print(f"\nComplete! Mode: {mode}")

# API - Application Programming Interface
def call_ollama_vision(image_path: str, model: str) -> str:
    """
    Call Ollama's vision model to analyze an image and suggest a filename.
    Returns: string of suggested filename (without extension)
    """
    # Read and encode the image
    with open(image_path, 'rb') as image_file:
        b64 = base64.b64encode(image_file.read()).decode('utf-8')
    # Prepare the API request
    url = "http://localhost:11434/api/generate"
    payload = {
        "model": model,
        "prompt": "Analyze this file's content and suggest a concise, 5-8 words, descriptive filename that captures its main content. Respond with ONLY the filename suggestion with NO additional text or explanation. Do NOT include a file extension.",
        "images": [b64],
        "stream": False
    }
    try:
        response = requests.post(url, json=payload)
        if response.status_code == 200:
            result = response.json()
            # Clean up the repsonse to get just the filename suggestion
            suggested_name = result['response'].strip()
            return suggested_name
        else:
            print(f"Error: API call failed with status code {response.status_code}")
            return None
    except Exception as e:
        print(f"Error calling Ollama API: {e}")
        return None

# API - Application Programming Interface 
def call_ollama_text(text, model: str) -> str:
    """
    Call Ollama's text model to analyze an text/code file and suggest a filename.
    Returns: string of suggested filename (without extension)
    """
    url = "http://localhost:11434/api/generate"
    payload = {
        "model": model,
        "prompt": "Analyze this file's content and suggest a concise, 5-8 word, descriptive filename that captures its main content. Respond with ONLY the filename suggestion with NO additional text or explanation. Do NOT include a file extension.\n\nFile Content:\n" + (text[:2000] if text else ''),
        "stream": False
    }
    try:
        response = requests.post(url, json=payload)
        if response.status_code == 200:
            result = response.json()
            # Clean up the repsonse to get just the filename suggestion
            suggested_name = result['response'].strip()
            return suggested_name
        else:
            print(f"Error: API call failed with status code {response.status_code}")
            return None
    except Exception as e:
        print(f"Error calling Ollama API: {e}")
        return None

# Extractor
def extract_pdf_image(file_path: str):
    # Try text extraction first
    try:
        reader = PdfReader(file_path)
        text = ""
        for i, page in enumerate(reader.pages[:3]):
            t = page.extract_text()
            if t:
                text += t
        if text.strip():
            return ("text", text[:2000])
    except Exception as e:
        print(f"PDF text extract error: {e}")
    # Fallback to image rendering of first page
    try:
        doc = fitz.open(file_path)
        page = doc[0]
        pix = page.get_pixmap(dpi=150)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp:
            pix.save(temp.name)
            temp_path = temp.name
        doc.close()
        return ("image", temp_path)
    except Exception as e:
        print(f"PDF image render error: {e}")
        return None

# Extractor
def extract_docx_content(file_path: str):
    try:
        doc = Document(file_path)
        text = '\n'.join(p.text for p in doc.paragraphs)
        return ("text", text[:2000])
    except Exception as e:
        print(f"DOCX extract error: {e}")
        return None

# Extractor
def extract_pptx_content(file_path: str):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return ("text", '\n'.join(text_runs)[:2000])
    except Exception as e:
        print(f"PPTX extract error: {e}")
        return None

# Extractor
def extract_video_frame(file_path: str):
    print("Video frame extraction not implemented yet")
    return None

# Extractor
def read_text_snippet(file_path: str):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return ("text", f.read(2000))
    except Exception as e:
        print(f"Text snippet error: {e}")
        return None

# Routing
def generate_filename_for_file(file_path: str, model: str):
    """
    Generate a new filename based on the file type and content.
    Returns: string of suggested filename (without extension)
    """
    # Get the file extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    # Image files handled as image
    if ext in ('.png', '.jpg', '.jpeg', '.webp', '.gif'):
        return call_ollama_vision(file_path, model)
    # PDF files
    elif ext == '.pdf':
        result = extract_pdf_image(file_path)
    # DOCX
    elif ext == '.docx':
        result = extract_docx_content(file_path)
    # PPTX
    elif ext == '.pptx':
        result = extract_pptx_content(file_path)
    # Video (not implemented)
    elif ext in ('.mp4', '.mov', '.avi'):
        result = extract_video_frame(file_path)
    # Text/Code files
    elif ext in ('.txt', '.md', '.csv', '.json', '.xml', '.py', '.java', '.c', '.cpp', '.html', '.css', '.js', '.ipynb'):
        result = read_text_snippet(file_path)
    else:
        return None
    if not result:
        return None
    media, data = result
    if media == "image":
        return call_ollama_vision(data, model)
    elif media == "text":
        return call_ollama_text(data, model)
    else:
        return None

# UX - User Experience
def apply_casing(text: str, case_style: str) -> str:
    """
    Apply the specified casing style to the text.
        text (str): Input text to transform
        case_style (str): One of 'snake', 'kebab', 'camel', 'pascal', 'lower', 'title'
    Returns: a string transformed text in the specified case style
    """
    # First, normalize the text by splitting into words
    # Convert to lowercase and split on spaces, dashes, underscores, or camel case boundaries
    words = [] # List to hold the extracted words
    current_word = "" # Temporary variable to build the current word
    for i, char in enumerate(text): # For each character in the text
        if char.isspace() or char in ['-', '_']: # Split on spaces, dashes, and underscores
            # Check that current_word is not empty
            if current_word:
                words.append(current_word.lower())
                current_word = ""
        # Split on camel case boundaries (uppercase letter directly following a lowercase letter)
        elif char.isupper() and i > 0 and text[i-1].islower():
            # Check that current_word is not empty
            if current_word:
                words.append(current_word.lower()) # Append the current word (lowercase) to words list
                current_word = char # Reset current_word
        else: # Otherwise, continue building the current word (because it's part of the same word)
            current_word += char
    if current_word: 
        words.append(current_word.lower()) # Append the last word if exists
    if not words:
        return ""
    # Apply the requested case style
    if case_style == 'snake':
        return '_'.join(words)
    elif case_style == 'kebab':
        return '-'.join(words)
    elif case_style == 'camel':
        return words[0].lower() + ''.join(word.capitalize() for word in words[1:])
    elif case_style == 'pascal':
        return ''.join(word.capitalize() for word in words)
    elif case_style == 'lower':
        return ' '.join(words)
    elif case_style == 'title':
        return ' '.join(word.capitalize() for word in words)
    else:
        return '_'.join(words) # Default to snake_case

def sanitize_filename(filename: str) -> str:
    """
    Sanitize the proposed filename to remove invalid characters and ensure it's safe for all operating systems.
    Returns: a string of the sanitized filename
    """
    # Define invalid characters (\\ is the escape char for \, allowing us to include '\' in the string)
    invalid_chars = '<>:"/\\|?*'
    # Replace invalid characters with underscores
    for char in invalid_chars:
        filename = filename.replace(char, '_')
    # Remove leading/trailing spaces and dots
    filename = filename.strip('. ')
    # Ensure the filename isn't empty
    if not filename:
        filename = 'unnamed'
    return filename

# Run main
if __name__ == "__main__":
    main()
